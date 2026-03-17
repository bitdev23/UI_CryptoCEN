"""Document processing: extract text from PDFs, DOCX, TXT, MD, CSV, PPTX and chunk for RAG."""
from typing import List, Tuple
import fitz  # pymupdf
import os
import csv
import io
from pathlib import Path
import logging

try:
    from docx import Document
    _HAS_DOCX = True
except Exception:
    _HAS_DOCX = False

try:
    from pptx import Presentation as PptxPresentation
    _HAS_PPTX = True
except Exception:
    _HAS_PPTX = False

logger = logging.getLogger("contentai.pdf_processor")

SUPPORTED_EXTENSIONS = ('.pdf', '.docx', '.txt', '.md', '.csv', '.pptx')


def extract_text_from_pdf(path: str) -> str:
    """Extract text from one PDF file; returns empty string on error."""
    try:
        doc = fitz.open(path)
    except Exception as e:
        logger.exception("Failed to open PDF: %s", path)
        return ""
    text_chunks = []
    try:
        for page in doc:
            text = page.get_text()
            if text:
                text_chunks.append(text)
    except Exception:
        logger.exception("Error while reading PDF pages: %s", path)
    finally:
        doc.close()
    return "\n".join(text_chunks)


def extract_text_from_docx(path: str) -> str:
    """Extract text from one DOCX file; returns empty string on error."""
    if not _HAS_DOCX:
        logger.warning("python-docx not installed; cannot parse DOCX: %s", path)
        return ""
    try:
        doc = Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
        return "\n".join(paragraphs)
    except Exception:
        logger.exception("Failed to extract DOCX: %s", path)
        return ""


def extract_text_from_txt(path: str) -> str:
    """Extract text from a plain text or markdown file."""
    try:
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()
    except Exception:
        logger.exception("Failed to extract TXT/MD: %s", path)
        return ""


def extract_text_from_csv(path: str) -> str:
    """Extract text from a CSV by converting rows to readable sentences."""
    try:
        rows = []
        with open(path, 'r', encoding='utf-8', errors='replace', newline='') as f:
            reader = csv.reader(f)
            headers = None
            for i, row in enumerate(reader):
                if i == 0:
                    headers = row
                    rows.append("Columns: " + ", ".join(row))
                else:
                    if headers:
                        parts = [f"{h}: {v}" for h, v in zip(headers, row) if v.strip()]
                        rows.append(" | ".join(parts))
                    else:
                        rows.append(" | ".join(row))
                if i > 2000:  # cap large CSVs
                    rows.append("... (truncated)")
                    break
        return "\n".join(rows)
    except Exception:
        logger.exception("Failed to extract CSV: %s", path)
        return ""


def extract_text_from_pptx(path: str) -> str:
    """Extract text from a PowerPoint presentation."""
    if not _HAS_PPTX:
        logger.warning("python-pptx not installed; cannot parse PPTX: %s", path)
        return ""
    try:
        prs = PptxPresentation(path)
        lines = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            lines.append(f"--- Slide {slide_num} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
        return "\n".join(lines)
    except Exception:
        logger.exception("Failed to extract PPTX: %s", path)
        return ""


def load_document(path: str) -> Tuple[str, str]:
    """Load one document file and return (source, text). Supports: PDF, DOCX, TXT, MD, CSV, PPTX."""
    """Load one document file and return (source, text). Supports: PDF, DOCX, TXT, MD, CSV, PPTX."""
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        text = extract_text_from_pdf(path)
    elif ext == ".docx":
        text = extract_text_from_docx(path)
    elif ext in (".txt", ".md"):
        text = extract_text_from_txt(path)
    elif ext == ".csv":
        text = extract_text_from_csv(path)
    elif ext == ".pptx":
        text = extract_text_from_pptx(path)
    else:
        logger.warning("Unsupported document type: %s", path)
        return (path, "")
    return (path, text)


def load_pdfs(folder: str = "data/pdfs") -> List[Tuple[str, str]]:
    """Load all supported documents in folder and return list of (filepath, text).
    Supported: .pdf, .docx, .txt, .md, .csv, .pptx
    """
    results = []
    p = Path(folder)
    if not p.exists():
        logger.warning("Document folder does not exist: %s", folder)
        return results
    for f in p.glob("**/*"):
        if f.suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue
        logger.info("Processing document: %s", f)
        try:
            source, text = load_document(str(f))
            if text:
                results.append((source, text))
            else:
                logger.warning("No text extracted from %s", f)
        except Exception:
            logger.exception("Failed to process document: %s", f)
    return results


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
    """Chunk text with overlap for RAG retrieval."""
    if chunk_size <= overlap:
        raise ValueError("chunk_size must be larger than overlap")
    if not text:
        return []

    out = []
    step = chunk_size - overlap
    start = 0
    L = len(text)
    while start < L:
        end = min(start + chunk_size, L)
        out.append(text[start:end])
        if end >= L:
            break
        start += step
    return out


if __name__ == "__main__":
    import dotenv
    dotenv.load_dotenv()
    logging.basicConfig(level=logging.DEBUG)
    pdfs = load_pdfs()
    print(f"Loaded {len(pdfs)} PDFs")
    if pdfs:
        chunks = chunk_text(pdfs[0][1])
        print(f"Example chunks: {len(chunks)}")
